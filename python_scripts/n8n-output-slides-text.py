import json
import os
from pptx import Presentation

# Get input data from items collection
items = $input.all()
results = []

# Process each item in the collection
for item in items:

    # Get file path from the current item
    file_path = item.json.get("file_path", "")
    
    if not file_path or not os.path.isfile(file_path):
        results.append({
            "success": False,
            "error": f"Invalid file path: {file_path}",
            "file_path": file_path
        })
        continue

def extract_text_from_shape(shape):
    """Extract text from a shape if it has a text frame"""
    if hasattr(shape, 'text'):
        return shape.text
    return ""

def extract_text_from_table(table):
    """Extract text from tables in the slide"""
    table_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            if cell.text:
                row_data.append(cell.text)
            else:
                row_data.append("")
        table_data.append(row_data)
    return table_data

def extract_slide_content(slide, slide_index):
    """Extract all text content from a slide"""
    slide_data = {
        "slide_number": slide_index + 1,
        "slide_id": str(slide.slide_id),
        "shapes_text": [],
        "tables": [],
        "notes": ""
    }
    
    # Extract text from shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                slide_data["shapes_text"].append(text)
        elif shape.has_table:
            table_data = extract_text_from_table(shape.table)
            if table_data:
                slide_data["tables"].append(table_data)
    
    # Extract speaker notes
    if hasattr(slide, 'notes_slide') and slide.notes_slide is not None:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            slide_data["notes"] = notes_text
    
    return slide_data

    try:
        # Load presentation
        presentation = Presentation(file_path)
        
        # Get presentation metadata
        presentation_data = {
            "title": os.path.basename(file_path),
            "slide_count": len(presentation.slides),
            "slides": []
        }
        
        # Extract content from each slide
        for i, slide in enumerate(presentation.slides):
            slide_data = extract_slide_content(slide, i)
            presentation_data["slides"].append(slide_data)
        
        # Output file path (in same directory as input file)
        base_name = os.path.splitext(file_path)[0]
        output_path = f"{base_name}_content.json"
        
        # Write to JSON file
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(presentation_data, f, ensure_ascii=False, indent=2)
        
        # Add result for this item
        results.append({
            "success": True,
            "file_path": file_path,
            "output_path": output_path,
            "slide_count": presentation_data["slide_count"],
            "data": presentation_data
        })

    except Exception as e:
        results.append({
            "success": False,
            "error": str(e),
            "file_path": file_path
        })

# Return all results
return results