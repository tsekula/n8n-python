# code for n8n Python Code node
# takes a file path and search and replace terms and writes out to disk an updated file

from pptx import Presentation
from io import BytesIO
import os

new_items = []
for item in items:
    try:
        # Get the file path from the item
        file_path = item.get('file_path', '')
        
        # Check if file path is provided and file exists
        if file_path and os.path.exists(file_path):
            # Get search and replace text (or use defaults if not provided)
            search_text = item.get('searchText', 'OLD TEXT')
            replace_text = item.get('replaceText', 'NEW TEXT')
            
            # Load the presentation from the file path
            prs = Presentation(file_path)
            
            # Process each slide
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        
                        for paragraph in text_frame.paragraphs:
                            if search_text in paragraph.text:
                                paragraph.text = paragraph.text.replace(search_text, replace_text)
            
            # Create output filename
            file_name = os.path.basename(file_path)
            output_path = os.path.join(os.path.dirname(file_path), f"modified_{file_name}")
            
            # Save the modified presentation
            prs.save(output_path)
            
            # Update item with results
            item['status'] = 'Text replaced successfully'
            item['output_path'] = output_path
            
        else:
            item['status'] = f"File not found: {file_path}"
            item['error'] = 'Invalid file path or file does not exist'
            
    except Exception as e:
        # Handle any errors
        item['status'] = 'Error processing PowerPoint file'
        item['error'] = str(e)
    
    new_items.append(item)

return new_items